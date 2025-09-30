import pytesseract
from pathlib import Path
from unidecode import unidecode
import spacy
import re
import phonenumbers
from rapidfuzz import process, fuzz
from spacy.matcher import PhraseMatcher
from docx import Document
from pdf2image import convert_from_path
import pdfplumber
from pptx import Presentation
from transformers import pipeline

# --- CONFIG ---
# Point pytesseract to the installed tesseract.exe location
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

SPACY_MODEL = "en_core_web_md"
resume_ner = pipeline("ner", model="dslim/bert-base-NER", aggregation_strategy="simple")
nlp = spacy.load(SPACY_MODEL)
edu_ner = pipeline("ner", model="microsoft/deberta-v3-base", aggregation_strategy="simple")

# Path to your skills vocabulary
SKILLS_FILE = Path(r"D:\TESTQ\resume_folder\paythonscript\skills.txt")

# --- Skills vocab load ---
def load_skills_vocab(skills_file: Path):
    lines = [unidecode(l.strip()) for l in skills_file.read_text(encoding="utf-8", errors="ignore").splitlines()]
    skills = {l.lower() for l in lines if l}
    matcher = PhraseMatcher(nlp.vocab, attr="LOWER")
    patterns = [nlp.make_doc(s) for s in skills]
    for i in range(0, len(patterns), 5000):
        matcher.add(f"SKILLS_{i//5000}", patterns[i:i+5000])
    return skills, matcher

skills_vocab, skills_matcher = load_skills_vocab(SKILLS_FILE)

# --- File readers ---
def read_pdf_text(path: Path) -> str:
    chunks = []
    with pdfplumber.open(str(path)) as pdf:
        for p in pdf.pages:
            chunks.append(p.extract_text() or "")
    return "\n".join(chunks)

def read_pdf_ocr(path: Path) -> str:
    text_parts = []
    try:
        images = convert_from_path(str(path), dpi=300)
        for img in images:
            txt = pytesseract.image_to_string(img, lang="eng")
            if txt.strip():
                text_parts.append(txt)
    except Exception as e:
        print(f"OCR failed for {path.name}: {e}")
    return "\n".join(text_parts)

def read_docx_text(path: Path) -> str:
    try:
        doc = Document(str(path))
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text)
        return "\n".join(texts)
    except Exception:
        return ""

def read_pptx_text(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""

def load_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext == ".pdf":
        txt = read_pdf_text(path)
        return txt if txt.strip() else read_pdf_ocr(path)
    if ext in [".docx", ".doc"]:
        return read_docx_text(path)
    if ext == ".pptx":
        return read_pptx_text(path)
    if ext in [".txt", ".rtf", ".md"]:
        return path.read_text(errors="ignore")
    return ""

# --- Core regexes & helpers ---
EMAIL_REGEX = re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b")
TITLE_HINTS = re.compile(
    r"\b(lead|senior|jr|junior|sr|principal|head|chief|manager|engineer|developer|analyst|scientist|consultant|designer|architect|administrator|specialist|intern|student|associate|officer|coordinator|executive|data|software|machine|ml|ai|devops|quality|qa|product|project|business)\b",
    re.I,
)
GPA_REGEX = re.compile(r"\b(?:CGPA|GPA)\s*[:=]?\s*(\d(?:\.\d{1,2})?)\s*(?:/|\bout of\b)?\s*(?:10|4)?", re.I)

EDU_KEYWORDS = re.compile(
    r"\b(education|educational background|academic|qualification|school|college|university|institute|academy|polytechnic|iit|iiit|nit)\b",
    re.I,
)
PROJECT_KEYWORDS = re.compile(r"\b(projects?|academic projects?|capstone|thesis)\b", re.I)
EXP_KEYWORDS = re.compile(r"\b(experience|work experience|employment|career|professional experience)\b", re.I)

ORG_HINTS = re.compile(
    r"\b(university|college|institute|school|academy|polytechnic|technological|technology|institute of|iit|iiit|nit|company|inc\.?|ltd\.?|llc|solutions|technologies|labs|systems)\b",
    re.I,
)

def extract_emails(text: str):
    # List of common email domain keywords
    email_domains = [
        "gmail.com", "yahoo.com", "outlook.com", "icloud.com", "hotmail.com",
        "rediffmail.com", "protonmail.com", "aol.com", "live.com", "msn.com",
        "yandex.com", "zoho.com", "mail.com", "gmx.com", "edu", "ac.in", "co.in"
    ]
    # Standard email regex
    email_regex = re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b")
    raw_emails = [e.rstrip('.,;:') for e in email_regex.findall(text)]

    # Also search for mailto: links
    mailto_regex = re.compile(r"mailto:([A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,})", re.I)
    mailto_emails = [m.group(1) for m in mailto_regex.finditer(text)]

    all_emails = raw_emails + mailto_emails

    # Prefer emails with known domains
    for email in all_emails:
        if any(domain in email.lower() for domain in email_domains):
            return email
    # If none match known domains, return the first found email if any
    if all_emails:
        return all_emails[0]
    return None

def extract_phone_numbers(text: str, default_region: str = None, max_candidates: int = 5):
    candidates = set()
    rough = re.findall(r"(?:\+?\d[\d\s().-]{6,}\d)", text)
    for r in rough[:200]:
        try:
            for match in phonenumbers.PhoneNumberMatcher(r, default_region or "IN"):
                num = match.number
                if phonenumbers.is_possible_number(num) and phonenumbers.is_valid_number(num):
                    candidates.add(phonenumbers.format_number(num, phonenumbers.PhoneNumberFormat.E164))
        except Exception:
            pass
    return list(candidates)[:max_candidates]

def extract_name(doc, text):
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    header_doc = nlp("\n".join(lines[:20]))
    persons = [ent for ent in header_doc.ents if ent.label_ == "PERSON"] or [ent for ent in doc.ents if ent.label_ == "PERSON"]
    return persons[0].text.strip() if persons else ""

def extract_designation(text):
    # List of popular designations
    popular_designations = [
        "software engineer", "data scientist", "project manager", "business analyst", "product manager",
        "developer", "designer", "consultant", "architect", "administrator", "specialist", "intern",
        "student", "associate", "officer", "coordinator", "executive", "qa engineer", "ml engineer",
        "ai engineer", "devops engineer", "quality analyst", "principal engineer", "senior engineer",
        "junior engineer", "lead engineer", "chief technology officer", "cto", "chief executive officer", "ceo"
    ]
    # Compile regex for popular designations
    designation_regex = re.compile(r"\b(" + "|".join(re.escape(d) for d in popular_designations) + r")\b", re.I)
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    for l in lines[:30]:
        m = designation_regex.search(l)
        if m:
            return m.group(1)
    return None

def extract_nationality(doc, text):
    m = re.search(r"(?:nationality|citizenship)\s*[:\-]\s*([A-Za-z ]{3,30})", text, re.I)
    return m.group(1).strip().title() if m else ""

def extract_skills(text: str, doc, skills_vocab, skills_matcher, fuzz_threshold=92):
    found = set()
    matches = skills_matcher(doc)
    for _, start, end in matches:
        found.add(doc[start:end].text)
    return sorted(found)

# ------- New: Education, GPA, Projects, Past Companies --------
def _find_section_blocks(lines, keyword_regex, window=80):
    """
    Return indexes of lines that likely start a section (Education/Projects/Experience),
    and the slice of lines belonging to that section until next section header.
    """
    headers = [i for i, l in enumerate(lines) if keyword_regex.search(l)]
    blocks = []
    for idx, h in enumerate(headers):
        start = h
        end = headers[idx + 1] if idx + 1 < len(headers) else min(len(lines), start + window)
        blocks.append((start, end))
    return blocks

def extract_education_and_gpa(text, doc=None):
    # Use NER to extract entities
    entities = edu_ner(text)
    institutions = set()
    results = []

    # Patterns for GPA
    gpa_regex = re.compile(r"\b(?:CGPA|GPA)\s*[:=]?\s*(\d(?:\.\d{1,2})?)\s*(?:/|\bout of\b)?\s*(?:10|4)?", re.I)

    # Collect institutes and GPA from NER
    for ent in entities:
        if ent["entity_group"] in ["ORG", "SCHOOL", "INSTITUTE", "UNIVERSITY"]:
            inst = ent["word"].strip()
            if inst.lower() not in institutions:
                institutions.add(inst.lower())
                # Search for GPA near the entity in text
                pattern = re.escape(inst) + r".{0,50}?" + gpa_regex.pattern
                m = re.search(pattern, text, re.I)
                gpa = m.group(1) if m else None
                results.append({"institution": inst, "gpa": gpa})

    # --- Fallback: keyword-based section parsing ---
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    blocks = _find_section_blocks(lines, EDU_KEYWORDS)
    for (s, e) in blocks:
        block = lines[s:e]
        for i, ln in enumerate(block):
            # Look for org-like education hints
            org_like = ORG_HINTS.search(ln) or any(ent.label_ == "ORG" and ent.text in ln for ent in nlp(ln).ents)
            if org_like:
                inst_clean = ln.strip()
                if inst_clean.lower() not in institutions:
                    institutions.add(inst_clean.lower())
                    # try GPA on same line, else next 2 lines
                    gpa = None
                    g = gpa_regex.search(ln)
                    if not g and i + 1 < len(block):
                        g = gpa_regex.search(block[i + 1])
                    if not g and i + 2 < len(block):
                        g = gpa_regex.search(block[i + 2])
                    if g:
                        gpa = g.group(1)
                    results.append({"institution": inst_clean, "gpa": gpa})
            else:
                # GPA lines that don't have institution; attach to last institution if any
                g = gpa_regex.search(ln)
                if g and results:
                    if results[-1]["gpa"] in (None, ""):
                        results[-1]["gpa"] = g.group(1)
    return results

def extract_projects(text):
    lines = [l.strip(" -•\t") for l in text.splitlines() if l.strip()]
    blocks = _find_section_blocks(lines, PROJECT_KEYWORDS)
    names = []

    def add(n):
        n = re.sub(r"\s{2,}", " ", n).strip()
        if n and n.lower() not in {x.lower() for x in names}:
            names.append(n)

    if not blocks:
        # Fallback: any line containing 'project'
        for ln in lines:
            if re.search(r"\bproject\b", ln, re.I) and len(ln) < 140:
                add(ln)
        return names

    for (s, e) in blocks:
        for ln in lines[s:e]:
            if re.search(r"\b(project|capstone|thesis)\b", ln, re.I):
                # Heuristic: keep first part before delimiter as the "name"
                name = re.split(r"[:\-–•]", ln, 1)[0]
                add(name)
    return names

def extract_past_companies(text: str, doc=None):
    """
    Extract past companies from resume text using spaCy ORG entities and regex,
    only including names with company suffixes and excluding programming languages/generic words.
    Deduplicate so only the most complete version of a company is kept.
    """
    companies_raw = []
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    blocks = _find_section_blocks(lines, EXP_KEYWORDS)

    # Company suffixes to look for (lowercase)
    company_suffixes = [
        "pvt ltd", "private limited", "inc", "llc", "ltd", "technologies", "solutions", "systems", "labs",
        "group", "corporation", "corp", "enterprises", "consulting", "industries", "services"
    ]
    # Programming languages to exclude (lowercase)
    programming_languages = [
        "python", "java", "c++", "c#", "javascript", "typescript", "go", "ruby", "php", "swift", "kotlin", "scala", "r", "matlab", "perl", "dart"
    ]
    # Education keywords to exclude
    edu_keywords = [
        "university", "college", "institute", "school", "academy", "iit", "nit", "polytechnic", "education"
    ]

    def has_company_suffix(name):
        name_l = name.lower()
        return any(re.search(r"\b" + re.escape(suf) + r"\b", name_l) for suf in company_suffixes)

    def is_valid_company(name):
        name_l = name.lower()
        # Must have a company suffix
        if not has_company_suffix(name):
            return False
        # Exclude if contains education keywords or programming languages
        if any(edu in name_l for edu in edu_keywords):
            return False
        if any(lang in name_l for lang in programming_languages):
            return False
        # Exclude if too short or generic
        if len(name.split()) < 2 or len(name) < 4:
            return False
        if name_l in ["project", "experience", "internship", "training"]:
            return False
        return True

    def extract_base_name(full_name):
        """
        Extract the base company name up to and including the first company suffix.
        """
        name = full_name
        for suf in company_suffixes:
            pattern = r"(.+?\b" + re.escape(suf) + r"\b)"
            m = re.search(pattern, name, re.I)
            if m:
                return m.group(1).strip().lower()
        return name.strip().lower()

    def add(c):
        c = re.sub(r"\s{2,}", " ", c).strip()
        if c and is_valid_company(c):
            companies_raw.append(c)

    # spaCy ORG entities from experience blocks
    if blocks:
        for (s, e) in blocks:
            seg = "\n".join(lines[s:e])
            seg_doc = nlp(seg)
            for ent in seg_doc.ents:
                if ent.label_ == "ORG":
                    add(ent.text)
            for ln in lines[s:e]:
                # Regex for company-like lines
                m = re.search(r"([A-Z][A-Za-z0-9&.,'()\- ]{2,})", ln)
                if m and has_company_suffix(m.group(1)):
                    add(m.group(1))
    else:
        gdoc = nlp(text)
        for ent in gdoc.ents:
            if ent.label_ == "ORG":
                add(ent.text)

    # Also extract from lines using suffixes
    for line in lines:
        for suf in company_suffixes:
            pattern = r"([A-Z][A-Za-z0-9&.,'()\- ]+?\s*" + re.escape(suf) + r"\b[^\n]*)"
            found = re.findall(pattern, line, re.I)
            for f in found:
                if is_valid_company(f):
                    companies_raw.append(f.strip())

    # Deduplicate: keep only the most complete entry for each base name
    base_to_full = {}
    for comp in companies_raw:
        base = extract_base_name(comp)
        # If this base is not seen or this entry is longer (more complete), keep it
        if base not in base_to_full or len(comp) > len(base_to_full[base]):
            base_to_full[base] = comp

    # Return sorted, deduplicated company names (title case)
    return sorted({v.strip().title() for v in base_to_full.values()})


# --- Public entry point ---
def extract_details_from_file(file_path: str):
    text = load_text(Path(file_path))
    if not text.strip():
        return {"error": "File is empty or unreadable"}

    text_ascii = unidecode(text)
    doc = nlp(text_ascii)

    # Existing fields
    base = {
        "name": extract_name(doc, text_ascii),
        "phone_number": "; ".join(extract_phone_numbers(text_ascii)),
        # "email": "; ".join(extract_emails(text_ascii)),
        "email": extract_emails(text_ascii),
        "designation": extract_designation(text_ascii),
        "skills": "; ".join(extract_skills(text_ascii, doc, skills_vocab, skills_matcher)),
        "nationality": extract_nationality(doc, text_ascii),
    }

    # New fields
    education = extract_education_and_gpa(text_ascii, doc)  
    projects = extract_projects(text_ascii)                 
    past_companies = extract_past_companies(text_ascii, doc)

    base.update({
        "education": education,
        "projects": projects,
        "past_companies": past_companies,
    })
    return base
